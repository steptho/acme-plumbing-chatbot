import streamlit as st
import os
import json
from datetime import datetime
from openai import OpenAI
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import base64
from io import BytesIO
from fpdf import FPDF
from fpdf import FPDF
import re

# -----------------------
# CONFIG
# -----------------------
st.set_page_config(page_title="Client Data Assistant", layout="wide")
client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

BASE_DIR = "client_data"
os.makedirs(BASE_DIR, exist_ok=True)

# -----------------------
# SESSION STATE INIT
# -----------------------
defaults = {
    "client_name": None,
    "messages": [],
    "current_chat_id": None,
    "uploaded_file": None
}

for key, value in defaults.items():
    if key not in st.session_state:
        st.session_state[key] = value

# -----------------------
# LOGIN
# -----------------------
def login():
    name = st.session_state.login_name.strip()
    if name:
        st.session_state.client_name = name
        st.session_state.messages = []
        st.session_state.current_chat_id = None

if not st.session_state.client_name:
    st.sidebar.title("🔒 Login")
    st.sidebar.text_input("Enter name/company", key="login_name", on_change=login)
    st.sidebar.button("Enter", on_click=login)
    st.stop()

CLIENT_DIR = os.path.join(BASE_DIR, st.session_state.client_name.replace(" ", "_"))
os.makedirs(CLIENT_DIR, exist_ok=True)

# -----------------------
# SAVE / LOAD CHAT
# -----------------------
def save_chat():
    if not st.session_state.current_chat_id:
        st.session_state.current_chat_id = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")

    path = os.path.join(CLIENT_DIR, f"{st.session_state.current_chat_id}.json")
    with open(path, "w") as f:
        json.dump(st.session_state.messages, f, indent=2)

def load_chat(file):
    with open(os.path.join(CLIENT_DIR, file), "r") as f:
        st.session_state.messages = json.load(f)
    st.session_state.current_chat_id = file.replace(".json","")

# -----------------------
# SIDEBAR
# -----------------------
st.sidebar.title(f"📂 {st.session_state.client_name}")

if st.sidebar.button("Logout"):
    st.session_state.clear()
    st.rerun()

st.sidebar.divider()

# Load existing chats
chat_files = sorted(
    [f for f in os.listdir(CLIENT_DIR) if f.endswith(".json")],
    reverse=True
)

for file in chat_files:
    if st.sidebar.button(file.replace(".json",""), key=file):
        load_chat(file)
        st.rerun()

# -----------------------
# FILE EXTRACTION
# -----------------------
def extract_text(file):
    text = ""
    try:
        if file.type.startswith("image"):
            base64_image = base64.b64encode(file.read()).decode("utf-8")
            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[{
                    "role":"user",
                    "content":[
                        {"type":"text","text":"Analyse this image professionally."},
                        {"type":"image_url","image_url":{"url":f"data:image/jpeg;base64,{base64_image}"}}
                    ]
                }]
            )
            return response.choices[0].message.content

        elif file.name.endswith(".pdf"):
            reader = PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() or ""
            return text

        elif file.name.endswith(".docx"):
            doc = Document(file)
            return "\n".join([p.text for p in doc.paragraphs])

        elif file.name.endswith(".pptx"):
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape,"text"):
                        text += shape.text + "\n"
            return text

        elif file.name.endswith(".xlsx"):
            df = pd.read_excel(file)
            return df.head(20).to_string()

        elif file.name.endswith(".csv"):
            df = pd.read_csv(file)
            return df.head(20).to_string()

    except Exception as e:
        return f"Error reading file: {e}"

    return "Unsupported file type."

# -----------------------
# PDF SAFE GENERATOR
# -----------------------
def generate_pdf(text):

    # Clean problematic characters
    text = text.replace("\r", "")
    text = re.sub(r"[^\x09\x0A\x20-\x7E\u00A0-\uFFFF]", "", text)

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=10)
    pdf.add_page()

    # Use Unicode-safe font
    font_path = "fonts/DejaVuSans.ttf"
    pdf.add_font("DejaVu", "", font_path, uni=True)
    pdf.set_font("DejaVu", size=10)
    # Break extremely long words
    def break_long_words(line, width=80):
        return [line[i:i+width] for i in range(0, len(line), width)]

    for line in text.split("\n"):
        safe_chunks = break_long_words(line)
        for chunk in safe_chunks:
            pdf.multi_cell(0, 5, chunk)

    pdf_bytes = BytesIO()
    pdf.output(pdf_bytes)
    pdf_bytes.seek(0)

    return pdf_bytes

# -----------------------
# MAIN UI
# -----------------------
st.title("💬 Document & Data Assistant")

# Upload file
uploaded = st.file_uploader("Upload PDF, Word, PPT, Excel, CSV, or Image")

if uploaded:
    st.session_state.uploaded_file = uploaded

    st.subheader("Live Preview")
    preview_text = extract_text(uploaded)
    st.text_area("File Content", preview_text[:3000], height=200)

    if st.button("🔍 Analyse File"):
        with st.spinner("Analysing..."):
            reply = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role":"system","content":"Provide a professional structured summary."},
                    {"role":"user","content":preview_text[:15000]}
                ]
            ).choices[0].message.content

            st.session_state.messages.append({"role":"assistant","content":reply})
            save_chat()
            st.success("Analysis Complete")
            st.rerun()

# -----------------------
# DISPLAY CHAT
# -----------------------
for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])

# Chat input
if prompt := st.chat_input("Ask a question about your document..."):
    st.session_state.messages.append({"role":"user","content":prompt})

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=st.session_state.messages
    )

    reply = response.choices[0].message.content
    st.session_state.messages.append({"role":"assistant","content":reply})
    save_chat()
    st.rerun()

# -----------------------
# DOWNLOAD PDF
# -----------------------
if st.session_state.messages:
    full_text = "\n\n".join(
        [f"{m['role'].upper()}: {m['content']}" for m in st.session_state.messages]
    )

    pdf_data = generate_pdf(full_text)

    st.download_button(
        "📄 Download Conversation as PDF",
        pdf_data,
        "chat_summary.pdf",
        "application/pdf"
    )



